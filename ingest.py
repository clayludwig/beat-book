"""
ingest.py
---------
Multi-format ingestion pipeline: raw bytes / URLs → normalized story dicts.

Two stages:
1. extract_text(filename, raw_bytes) → plain text/markdown
   Format dispatcher. markitdown handles docx/pdf/html/pptx/xlsx/rtf.
   stdlib handles txt/md/json. Unknown extensions get a utf-8 best-effort.

2. normalize(text, source_label, anthropic_key) → list[Story]
   Claude Sonnet 4.6 tool-use call with a strict schema. Documents that fit
   in one window get a single call; larger documents are split into
   overlapping windows, processed concurrently, and deduplicated on merge.
   The LLM only infers metadata (title/date/author) and returns verbatim
   markers for each story's body — the body itself is sliced from the
   original text, never LLM-rewritten.

A wrapper, ingest_source(), runs both stages and packages the result for
the /ingest route.
"""

from __future__ import annotations

import concurrent.futures
import html
import ipaddress
import json
import logging
import re
import socket
import tempfile
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Optional, Tuple
import threading
from urllib.parse import urlparse

import anthropic
import httpx

from claude_client import (
    ANTHROPIC_SEMAPHORE,
    CHAT_MODEL,
    RATE_LIMIT_MAX_RETRIES,
    chat_client,
    rate_limit_pause,
)

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────────────────────────────────────

MAX_FILE_BYTES = 25 * 1024 * 1024            # 25 MB per file
URL_FETCH_TIMEOUT = 15.0                     # seconds
URL_USER_AGENT = "BeatBookBuilder/1.0 (+https://github.com/clayludwig/beat-book)"
NORMALIZE_MODEL = "claude-haiku-4-5-20251001"
# OCR settings for scanned PDFs (rendered via PyMuPDF, transcribed via Haiku vision).
OCR_DPI = 150
OCR_PAGES_PER_BATCH = 4   # images per Haiku vision call
OCR_MAX_PAGES = 100        # cap for very large scanned PDFs
# Max tokens per LLM call. Sized to fit marker data for up to ~30 stories
# per chunk.
NORMALIZE_MAX_TOKENS = 4096
# Number of attempts per chunk when the model returns no tool call. With
# tool_choice forced this should be rare, but one retry covers transients.
NORMALIZE_MAX_ATTEMPTS = 2

# Chunked normalization: documents over WINDOW_SIZE get split into overlapping
# windows and processed concurrently. WINDOW_OVERLAP must be larger than any
# expected single-story body so a story spanning a boundary is fully contained
# in one of the two adjacent windows.
WINDOW_SIZE = 100_000           # primary chunk size in chars
WINDOW_OVERLAP = 15_000          # chars of overlap between adjacent chunks
# Soft safety cap on LLM calls per document. Sized to comfortably cover
# the largest input the 25 MB raw-file cap could produce after extraction
# (≈50 MB of HTML-stripped text at ~100 KB per chunk).
MAX_CHUNKS = 500
NORMALIZE_CONCURRENCY = 4
# Record separator emitted by _extract_json for top-level JSON lists.
# When present in the extracted text, _make_chunks packs whole records
# into each chunk so a story body is never split across two chunks.
RECORD_SEPARATOR = "\n\n---\n\n"

# Extensions markitdown converts to readable markdown.
# Extensions we read directly as utf-8 text.
_TEXT_EXTS = {".txt", ".md", ".markdown", ".log", ".csv"}


# ─────────────────────────────────────────────────────────────────────────────
# TYPES
# ─────────────────────────────────────────────────────────────────────────────


class IngestError(Exception):
    """Raised when a single source can't be ingested. Surfaced to the user."""


CONTENT_TYPES = frozenset({
    "article", "document", "dataset", "report",
    "transcript", "press_release", "post", "other",
})


@dataclass
class Story:
    """A single normalized content entry. Output of stage 2."""
    title: str
    content: str
    date: str = ""
    author: str = ""
    organization: str = ""
    link: str = ""
    content_type: str = "article"   # one of CONTENT_TYPES
    metadata: dict = field(default_factory=dict)  # type-specific fields
    confidence: str = "medium"      # "high" | "medium" | "low"
    reasoning: str = ""             # one-sentence justification for the preview UI

    def to_pipeline_dict(self) -> dict:
        """Shape consumed by pipeline.py / agent.py / citation_matcher.py."""
        out = {"title": self.title, "content": self.content}
        if self.date:
            out["date"] = self.date
        if self.author:
            out["author"] = self.author
        if self.organization:
            out["organization"] = self.organization
        if self.link:
            out["link"] = self.link
        if self.content_type:
            out["content_type"] = self.content_type
        if self.metadata:
            out["metadata"] = self.metadata
        return out

    def to_preview_dict(self) -> dict:
        """Shape sent to the preview UI (includes metadata fields)."""
        return {
            "title": self.title,
            "content": self.content,
            "date": self.date,
            "author": self.author,
            "organization": self.organization,
            "link": self.link,
            "content_type": self.content_type,
            "metadata": self.metadata,
            "confidence": self.confidence,
            "reasoning": self.reasoning,
        }


@dataclass
class IngestedSource:
    """Result of ingesting one file or one URL."""
    source_label: str            # filename or URL
    kind: str                    # "file" | "url"
    stories: list[Story] = field(default_factory=list)
    excluded: bool = False
    skip_reason: str = ""
    extract_error: str = ""
    char_count: int = 0
    truncated: bool = False

    def to_preview_dict(self) -> dict:
        return {
            "source_label": self.source_label,
            "kind": self.kind,
            "stories": [s.to_preview_dict() for s in self.stories],
            "excluded": self.excluded,
            "skip_reason": self.skip_reason,
            "extract_error": self.extract_error,
            "char_count": self.char_count,
            "truncated": self.truncated,
        }


# ─────────────────────────────────────────────────────────────────────────────
# STRUCTURED JSON FAST PATH
# ─────────────────────────────────────────────────────────────────────────────
# When a JSON file is a recognisable list of story-like objects (RSS feeds,
# news API responses, the chicago-public-media daily exports) the metadata
# is already explicit — no LLM call needed. We detect the structure, map
# fields directly, and skip normalization entirely.

_STORY_CONTENT_KEYS = ("summary", "content", "body", "text", "description",
                        "content_html", "content:encoded")
_STORY_DATE_KEYS    = ("published", "date", "pubDate", "pub_date",
                        "created_at", "updated_at", "timestamp")
_STORY_AUTHOR_KEYS  = ("author", "byline", "creator", "dc:creator")
_STORY_LINK_KEYS    = ("link", "url", "href", "guid")
_STORY_LIST_KEYS    = ("entries", "items", "stories", "articles",
                        "results", "data", "posts", "feed")
_MIN_CONTENT_CHARS  = 50


def _extract_story_list(data) -> Optional[list]:
    if isinstance(data, list):
        return data
    if isinstance(data, dict):
        for key in _STORY_LIST_KEYS:
            v = data.get(key)
            if isinstance(v, list) and v:
                return v
    return None


def _rendered_field(value) -> str:
    """Return a string from a raw value or a WP-style {rendered: ...} object."""
    if isinstance(value, dict):
        rendered = value.get("rendered")
        return str(rendered).strip() if rendered else ""
    if isinstance(value, str):
        return value.strip()
    return ""


def _looks_like_story_list(items: list) -> bool:
    sample = [x for x in items[:10] if isinstance(x, dict)]
    if not sample:
        return False
    has_title = sum(1 for x in sample if _rendered_field(x.get("title")))
    has_content = sum(
        1 for x in sample
        if any(_rendered_field(x.get(k)) for k in _STORY_CONTENT_KEYS)
    )
    return has_title >= len(sample) * 0.8 and has_content >= len(sample) * 0.5


def _map_json_item(item: dict, link_hint: str) -> Optional["Story"]:
    title = _rendered_field(item.get("title"))
    if not title:
        return None

    content_raw = ""
    for k in _STORY_CONTENT_KEYS:
        content_raw = _rendered_field(item.get(k))
        if content_raw:
            break
    content = _clean_inline_html(content_raw).strip()
    if len(content) < _MIN_CONTENT_CHARS:
        return None

    date = ""
    for k in _STORY_DATE_KEYS:
        v = item.get(k)
        if v:
            m = re.match(r"(\d{4}-\d{2}-\d{2})", str(v).strip())
            if m:
                date = m.group(1)
                break

    author = ""
    for k in _STORY_AUTHOR_KEYS:
        v = item.get(k)
        if v:
            author = str(v.get("name") if isinstance(v, dict) else v).strip()
            if author:
                break

    link = link_hint
    for k in _STORY_LINK_KEYS:
        v = item.get(k)
        if isinstance(v, dict):
            v = v.get("rendered")
        if v and isinstance(v, str) and v.startswith("http"):
            link = v.strip()
            break

    return Story(
        title=title, content=content, date=date, author=author, link=link,
        content_type="article",
        confidence="high",
        reasoning="Mapped directly from structured JSON fields.",
    )


def _fast_json_stories(
    raw: bytes, source_label: str, link_hint: str = ""
) -> Optional[list["Story"]]:
    """Return stories extracted from a structured JSON file without an LLM
    call, or None if the structure doesn't match and normal normalization
    should be used."""
    try:
        data = json.loads(raw.decode("utf-8", errors="replace"))
    except (json.JSONDecodeError, ValueError):
        return None

    items = _extract_story_list(data)
    if not items or not _looks_like_story_list(items):
        return None

    stories = [
        s for item in items
        if isinstance(item, dict)
        for s in [_map_json_item(item, link_hint)]
        if s is not None
    ]
    return stories if stories else None


# ─────────────────────────────────────────────────────────────────────────────
# STAGE 1 — EXTRACT
# ─────────────────────────────────────────────────────────────────────────────


def _ext_of(filename: str) -> str:
    return Path(filename).suffix.lower()


def _decode_text(raw: bytes) -> str:
    """Best-effort utf-8 decode with replacement for stray bytes."""
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("utf-8", errors="replace")


def _clean_inline_html(s: str) -> str:
    """Unescape HTML entities and strip tags, preserving paragraph breaks.

    Applied to bodies pulled out of JSON/RSS records where the value commonly
    contains raw HTML. Keeps the text in a shape the LLM can read directly,
    so the markers it returns will match the extracted text.
    """
    if not s:
        return s
    # Decode entities twice in case of double-encoding (&amp;lt; → &lt; → <).
    text = html.unescape(html.unescape(s))
    # Block tags → paragraph breaks before stripping the rest.
    text = re.sub(r"(?i)<\s*br\s*/?\s*>", "\n", text)
    text = re.sub(
        r"(?i)</\s*(p|div|li|h[1-6]|blockquote|tr|article|section)\s*>",
        "\n\n", text,
    )
    text = re.sub(
        r"(?i)<\s*(p|div|li|h[1-6]|blockquote|tr|article|section)(\s[^>]*)?>",
        "\n\n", text,
    )
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_pdf(raw: bytes) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(stream=raw, filetype="pdf")
    n_pages = len(doc)
    pages = [page.get_text().strip() for page in doc if page.get_text().strip()]
    doc.close()
    if not pages and n_pages > 0:
        raise IngestError("__SCANNED_PDF__")
    return "\n\n".join(pages)


def _render_page_png(page) -> bytes:
    import fitz
    mat = fitz.Matrix(OCR_DPI / 72, OCR_DPI / 72)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    return pix.tobytes("png")


def _ocr_pdf(raw: bytes, anthropic_key: str, source_label: str) -> str:
    """OCR a scanned PDF using Claude Haiku vision.

    Renders pages to PNG at OCR_DPI, batches them into groups of
    OCR_PAGES_PER_BATCH, sends each batch to Haiku for transcription,
    and returns the concatenated text. Caps at OCR_MAX_PAGES.
    """
    import base64
    import fitz

    doc = fitz.open(stream=raw, filetype="pdf")
    n_pages = len(doc)
    if n_pages == 0:
        doc.close()
        raise IngestError(f"{source_label}: PDF has no pages.")

    truncated = n_pages > OCR_MAX_PAGES
    page_indices = list(range(min(n_pages, OCR_MAX_PAGES)))

    page_pngs: list[bytes] = []
    for i in page_indices:
        page_pngs.append(_render_page_png(doc[i]))
    doc.close()

    client = chat_client(anthropic_key)

    batches = [page_pngs[i:i + OCR_PAGES_PER_BATCH]
               for i in range(0, len(page_pngs), OCR_PAGES_PER_BATCH)]
    batch_page_nums = [page_indices[i:i + OCR_PAGES_PER_BATCH]
                       for i in range(0, len(page_indices), OCR_PAGES_PER_BATCH)]

    def _ocr_batch(pngs: list[bytes], nums: list[int]) -> str:
        content: list[dict] = []
        for png in pngs:
            b64 = base64.standard_b64encode(png).decode()
            content.append({
                "type": "image",
                "source": {"type": "base64", "media_type": "image/png", "data": b64},
            })
        label = (f"pages {nums[0]+1}–{nums[-1]+1}" if len(nums) > 1
                 else f"page {nums[0]+1}")
        content.append({
            "type": "text",
            "text": (
                f"These are {len(pngs)} consecutive pages ({label}) from a scanned PDF. "
                "Transcribe ALL text exactly as it appears, in reading order. "
                "Preserve paragraph structure with blank lines between paragraphs. "
                "Output only the transcribed text — no commentary, no page labels."
            ),
        })
        for rl_attempt in range(RATE_LIMIT_MAX_RETRIES + 1):
            try:
                with ANTHROPIC_SEMAPHORE:
                    resp = client.messages.create(
                        model=NORMALIZE_MODEL,
                        max_tokens=4096,
                        messages=[{"role": "user", "content": content}],
                    )
                break
            except anthropic.RateLimitError as e:
                if rl_attempt >= RATE_LIMIT_MAX_RETRIES:
                    raise
                pause = rate_limit_pause(rl_attempt, e)
                logger.warning("OCR rate limited; waiting %.0fs (attempt %d/%d).",
                               pause, rl_attempt + 1, RATE_LIMIT_MAX_RETRIES)
                time.sleep(pause)
        return "".join(
            b.text for b in resp.content if getattr(b, "type", None) == "text"
        )

    batch_results: list[str | None] = [None] * len(batches)
    with concurrent.futures.ThreadPoolExecutor(max_workers=NORMALIZE_CONCURRENCY) as ex:
        futures = {
            ex.submit(_ocr_batch, batches[i], batch_page_nums[i]): i
            for i in range(len(batches))
        }
        for fut in concurrent.futures.as_completed(futures):
            i = futures[fut]
            try:
                batch_results[i] = fut.result()
            except Exception as e:
                logger.warning(
                    "OCR batch %d/%d failed for %s: %s",
                    i + 1, len(batches), source_label, e,
                )
                batch_results[i] = ""

    combined = "\n\n".join(t for t in batch_results if t).strip()
    if not combined:
        raise IngestError(
            f"{source_label}: OCR produced no text. The PDF may be too degraded to read."
        )
    if truncated:
        combined += (
            f"\n\n[Note: This PDF has {n_pages} pages. "
            f"Only the first {OCR_MAX_PAGES} were processed by OCR.]"
        )
    return combined


# Filename keyword → content type, checked before falling back to "document".
_FILENAME_TYPE_HINTS: list[tuple[tuple[str, ...], str]] = [
    (("minutes", "meeting_minutes"),                     "transcript"),
    (("newsletter", "bulletin"),                         "report"),
    (("report", "action_report", "disciplinary_action"), "report"),
    (("agenda",),                                        "document"),
    (("license", "licensure", "license_actions"),        "document"),
    (("order", "suspension", "revocation", "reprimand",
      "surrender", "reinstatement", "consent_order",
      "charges", "complaint", "final_order"),             "document"),
]


def _infer_content_type_from_filename(filename: str) -> str:
    """Return a content_type hint derived from the filename, or 'document'."""
    name = Path(filename).stem.lower()
    for keywords, ctype in _FILENAME_TYPE_HINTS:
        if any(kw in name for kw in keywords):
            return ctype
    return "document"



def _extract_docx(raw: bytes) -> str:
    import io
    import docx
    doc = docx.Document(io.BytesIO(raw))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    # Tables
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _extract_pptx(raw: bytes) -> str:
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_xlsx(raw: bytes) -> str:
    import io
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    sheets = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def _extract_xls(raw: bytes) -> str:
    import io
    import xlrd
    wb = xlrd.open_workbook(file_contents=raw)
    sheets = []
    for sheet in wb.sheets():
        rows = []
        for rx in range(sheet.nrows):
            cells = [str(sheet.cell_value(rx, cx)).strip()
                     for cx in range(sheet.ncols)
                     if str(sheet.cell_value(rx, cx)).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def _extract_html(raw: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n").strip()


def _extract_rtf(raw: bytes) -> str:
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(_decode_text(raw)).strip()


def _extract_epub(raw: bytes) -> str:
    import io
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    book = epub.read_epub(io.BytesIO(raw))
    parts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        text = soup.get_text(separator="\n").strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _indent(text: str, by: str = "  ") -> str:
    return "\n".join((by + ln) if ln else ln for ln in text.split("\n"))


def _render_value(v) -> str:
    """Render any JSON value as readable text. No field names are
    privileged — every key is labeled, every value rendered recursively.
    Strings get HTML stripped so Stage 2's verbatim-marker resolver can
    find them in the source."""
    if v is None or v == "" or v == [] or v == {}:
        return ""
    if isinstance(v, bool):
        return "true" if v else "false"
    if isinstance(v, str):
        return _clean_inline_html(v)
    if isinstance(v, (int, float)):
        return str(v)
    if isinstance(v, list):
        rendered = [s for s in (_render_value(x) for x in v) if s]
        if not rendered:
            return ""
        if all("\n" not in s and len(s) < 60 for s in rendered):
            return ", ".join(rendered)
        return "\n".join(f"- {_indent(s).lstrip()}" for s in rendered)
    if isinstance(v, dict):
        parts = []
        for k, sub in v.items():
            sub_text = _render_value(sub)
            if not sub_text:
                continue
            if "\n" in sub_text:
                parts.append(f"{k}:\n{_indent(sub_text)}")
            else:
                parts.append(f"{k}: {sub_text}")
        return "\n".join(parts)
    return str(v)


def _extract_json(raw: bytes) -> str:
    """Render JSON as readable text the LLM can scan for stories.

    Schema-agnostic — no field names are special. A top-level list is
    split with `---` separators so each item becomes a candidate story
    boundary; anything else is rendered as one document. HTML embedded
    in string values is stripped so Stage 2's verbatim markers resolve
    against the rendered text.
    """
    text = _decode_text(raw)
    try:
        data = json.loads(text)
    except json.JSONDecodeError:
        # Not valid JSON — return raw text; the LLM can still try.
        return text

    if isinstance(data, list):
        rendered = [r for r in (_render_value(item) for item in data) if r]
        return "\n\n---\n\n".join(rendered)
    return _render_value(data)


def extract_text(filename: str, raw: bytes) -> str:
    """Stage 1. Dispatch on extension to produce clean text/markdown."""
    if len(raw) > MAX_FILE_BYTES:
        raise IngestError(
            f"{filename}: file is {len(raw) / 1_048_576:.1f} MB; the limit is "
            f"{MAX_FILE_BYTES / 1_048_576:.0f} MB."
        )

    ext = _ext_of(filename)

    _EXTRACTORS = {
        ".json":  lambda: _extract_json(raw),
        ".pdf":   lambda: _extract_pdf(raw),
        ".docx":  lambda: _extract_docx(raw),
        ".doc":   lambda: _extract_docx(raw),
        ".pptx":  lambda: _extract_pptx(raw),
        ".ppt":   lambda: _extract_pptx(raw),
        ".xlsx":  lambda: _extract_xlsx(raw),
        ".xls":   lambda: _extract_xls(raw),
        ".html":  lambda: _extract_html(raw),
        ".htm":   lambda: _extract_html(raw),
        ".rtf":   lambda: _extract_rtf(raw),
        ".epub":  lambda: _extract_epub(raw),
    }

    if ext in _TEXT_EXTS:
        return _decode_text(raw).strip()

    if ext in _EXTRACTORS:
        try:
            return _EXTRACTORS[ext]()
        except Exception as e:
            raise IngestError(
                f"{filename}: failed to extract ({ext}) — {type(e).__name__}: {e}"
            ) from e

    # Unknown extension — try utf-8 decode, reject if it looks binary.
    decoded = _decode_text(raw).strip()
    if decoded and "\x00" not in decoded[:1000]:
        return decoded
    raise IngestError(
        f"{filename}: unsupported file type {ext!r}; could not extract text."
    )


# ─────────────────────────────────────────────────────────────────────────────
# URL FETCHING (with SSRF protection)
# ─────────────────────────────────────────────────────────────────────────────


def _is_blocked_ip(host: str) -> bool:
    """Resolve hostname; reject loopback / private / link-local / multicast."""
    try:
        # getaddrinfo returns all A/AAAA records — check every one.
        infos = socket.getaddrinfo(host, None)
    except socket.gaierror:
        return True  # unresolvable → block

    for info in infos:
        ip_str = info[4][0]
        try:
            ip = ipaddress.ip_address(ip_str)
        except ValueError:
            continue
        if (
            ip.is_private
            or ip.is_loopback
            or ip.is_link_local
            or ip.is_multicast
            or ip.is_reserved
            or ip.is_unspecified
        ):
            return True
    return False


def fetch_url(url: str) -> Tuple[str, bytes]:
    """Fetch a URL, return (suggested_filename, raw_bytes).

    Refuses non-http(s) schemes and private/loopback addresses.
    Honors MAX_FILE_BYTES.
    """
    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        raise IngestError(f"{url}: only http and https URLs are accepted.")
    if not parsed.hostname:
        raise IngestError(f"{url}: URL is missing a hostname.")
    if _is_blocked_ip(parsed.hostname):
        raise IngestError(f"{url}: URL resolves to a private or unreachable address.")

    headers = {"User-Agent": URL_USER_AGENT, "Accept": "*/*"}
    try:
        with httpx.Client(
            timeout=URL_FETCH_TIMEOUT,
            follow_redirects=True,
            headers=headers,
        ) as client:
            resp = client.get(url)
    except httpx.HTTPError as e:
        raise IngestError(f"{url}: fetch failed — {type(e).__name__}: {e}") from e

    if resp.status_code >= 400:
        raise IngestError(f"{url}: server returned HTTP {resp.status_code}.")

    body = resp.content
    if len(body) > MAX_FILE_BYTES:
        raise IngestError(
            f"{url}: response is {len(body) / 1_048_576:.1f} MB; the limit is "
            f"{MAX_FILE_BYTES / 1_048_576:.0f} MB."
        )

    # Decide a filename for extension dispatch.
    content_type = (resp.headers.get("content-type") or "").split(";")[0].strip().lower()
    type_to_ext = {
        "text/html": ".html",
        "application/xhtml+xml": ".html",
        "text/plain": ".txt",
        "text/markdown": ".md",
        "application/json": ".json",
        "application/pdf": ".pdf",
        "application/msword": ".doc",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
        "application/rtf": ".rtf",
        "text/rtf": ".rtf",
    }
    suggested_ext = type_to_ext.get(content_type, "")
    path_part = Path(parsed.path).name or parsed.hostname
    if suggested_ext and not path_part.lower().endswith(suggested_ext):
        path_part = f"{path_part}{suggested_ext}" if path_part else f"page{suggested_ext}"
    elif not _ext_of(path_part):
        path_part = f"{path_part}.html"  # default for unknown content-type

    return path_part, body


# ─────────────────────────────────────────────────────────────────────────────
# STAGE 2 — NORMALIZE (LLM)
# ─────────────────────────────────────────────────────────────────────────────


_NORMALIZE_TOOL = {
    "name": "register_stories",
    "description": (
        "Register the content entries found in the document. "
        "Return is_news_content=false only when the document has no journalistic value "
        "(empty file, binary garbage, personal shopping list) — "
        "in that case return an empty stories list and explain in skip_reason."
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "is_news_content": {
                "type": "boolean",
                "description": (
                    "True if the document contains anything useful for journalism research — "
                    "news articles, data, reports, records, legal documents, spreadsheets, etc. "
                    "False ONLY for empty files, binary garbage, or content with zero journalistic value."
                ),
            },
            "skip_reason": {
                "type": "string",
                "description": "If is_news_content is false, one short sentence explaining why. Empty otherwise.",
            },
            "stories": {
                "type": "array",
                "description": (
                    "One entry per distinct unit of content. "
                    "For news: one entry per article. "
                    "For datasets: one entry per logical table (NOT one per row). "
                    "For long reports: one entry per major section where natural. "
                    "Leave empty if is_news_content is false."
                ),
                "items": {
                    "type": "object",
                    "properties": {
                        "content_type": {
                            "type": "string",
                            "enum": ["article", "document", "dataset", "report",
                                     "transcript", "press_release", "post", "other"],
                            "description": (
                                "Identify the content type FIRST — it determines which metadata fields to populate.\n"
                                "article = news story, op-ed, feature, blog post\n"
                                "document = government filing, legal order, license record, FOIA doc, policy doc\n"
                                "dataset = structured/tabular data: CSV, spreadsheet, database export\n"
                                "report = research report, audit, study, white paper, annual report, board newsletter\n"
                                "transcript = interview, hearing, deposition, meeting minutes\n"
                                "press_release = official statement or announcement\n"
                                "post = social media post, forum thread\n"
                                "other = none of the above"
                            ),
                        },
                        "title": {
                            "type": "string",
                            "description": (
                                "A clear, specific title.\n"
                                "article: the headline.\n"
                                "document (legal order): the document's own title, e.g. 'Order of Summary Suspension'.\n"
                                "document (license): 'License Record — [Name] ([State])'.\n"
                                "dataset: describe the data, e.g. 'Kentucky Disciplinary Records 2020–2024'.\n"
                                "transcript: '[Body Name] Meeting Minutes — [Date]'.\n"
                                "report: the report's own title.\n"
                                "If no title is present, write a 6-10 word descriptive title."
                            ),
                        },
                        "date": {
                            "type": "string",
                            "description": (
                                "Most relevant date as YYYY-MM-DD.\n"
                                "article: publication date.\n"
                                "document (legal): order/decision/effective date.\n"
                                "transcript: meeting date.\n"
                                "report: release or cover date.\n"
                                "Empty string if not present."
                            ),
                        },
                        "author": {
                            "type": "string",
                            "description": (
                                "An individual person — not an institution.\n"
                                "article: byline author.\n"
                                "document: signing official or hearing officer if named.\n"
                                "Leave empty for institutional documents; use 'organization' instead."
                            ),
                        },
                        "organization": {
                            "type": "string",
                            "description": (
                                "The institution associated with this entry.\n"
                                "article: publication name (e.g. 'The New York Times').\n"
                                "document: issuing board, court, or agency (e.g. 'Idaho Board of Medicine').\n"
                                "dataset: organization that produced or owns the data.\n"
                                "transcript: committee or board that met.\n"
                                "report: issuing organization.\n"
                                "Empty if not identifiable."
                            ),
                        },
                        "link": {
                            "type": "string",
                            "description": "Source URL if present in the document. Empty string otherwise.",
                        },
                        "metadata": {
                            "type": "object",
                            "description": (
                                "Type-specific fields. Only include keys whose values are explicitly present. Use snake_case.\n\n"
                                "article → publication_section, word_count\n"
                                "document (legal order/action) → docket_number, action_type (e.g. Summary Suspension / Revocation / Consent Order / Public Reprimand / Reinstatement), subject_name (person or org the action concerns), jurisdiction (state/federal)\n"
                                "document (license record) → license_number, license_status (Active/Suspended/Expired/Revoked), state, expiry_date\n"
                                "dataset → source_organization, date_range (e.g. 2020-2024), fields (comma-separated column names), row_count\n"
                                "report → issuing_organization, report_number\n"
                                "transcript → body_name (full name of committee/board), meeting_format (in-person/virtual/hybrid), key_attendees\n"
                                "press_release → contact_name, contact_email\n"
                                "post → platform (Twitter/X/Facebook/etc.), handle"
                            ),
                            "additionalProperties": True,
                        },
                        "body_starts_with": {
                            "type": "string",
                            "description": (
                                "First 30-80 characters of this entry's body, COPIED VERBATIM from the document. "
                                "Must appear exactly — used for substring search. "
                                "Start after the title/header/byline."
                            ),
                        },
                        "body_ends_with": {
                            "type": "string",
                            "description": (
                                "Last 30-80 characters of this entry's body, COPIED VERBATIM from the document. "
                                "Must appear exactly — used for substring search."
                            ),
                        },
                        "confidence": {
                            "type": "string",
                            "enum": ["high", "medium", "low"],
                            "description": "high = metadata explicit in text. medium = reasonably inferred. low = a guess.",
                        },
                        "reasoning": {
                            "type": "string",
                            "description": "One sentence: what this entry is and how you identified it.",
                        },
                    },
                    "required": [
                        "content_type", "title", "date", "author", "organization",
                        "link", "metadata", "body_starts_with", "body_ends_with",
                        "confidence", "reasoning",
                    ],
                },
            },
        },
        "required": ["is_news_content", "skip_reason", "stories"],
    },
}


_NORMALIZE_SYSTEM = (
    "You are a document parser for the Beat Book Builder, a tool that helps "
    "journalists turn source materials into a reporting guide. "
    "Source materials include news articles, datasets, government records, "
    "legal orders, licenses, board minutes, reports, spreadsheets, "
    "CSV files, JSON feeds, and any other journalism-relevant content.\n\n"
    "STEP 1 — Identify what the document is.\n"
    "Read enough to classify it: is this a news article? A legal order? "
    "A license record? Meeting minutes? A dataset? A report?\n\n"
    "STEP 2 — Extract entries appropriate to that type.\n"
    "- One document may contain ZERO, ONE, or MANY entries.\n"
    "- Split on natural boundaries: headings, byline blocks, record separators, sections.\n"
    "- article: one entry per story. Title = headline. Body = article text.\n"
    "- document (legal order/action): one entry per order. Title = the document's own title "
    "(e.g. 'Order of Summary Suspension'). Date = order/effective date. "
    "Metadata: docket_number, action_type, subject_name, jurisdiction.\n"
    "- document (license): one entry per license record. "
    "Metadata: license_number, license_status, state, expiry_date.\n"
    "- dataset/tabular: one entry per logical table — NOT one per row. "
    "Title describes the data. Body = full table text. "
    "Metadata: fields, date_range, source_organization.\n"
    "- report: one entry per document (or per major section for long reports). "
    "Metadata: issuing_organization, report_number.\n"
    "- transcript: one entry per meeting/session. Title = body + date. "
    "Metadata: body_name, meeting_format, key_attendees.\n\n"
    "STEP 3 — Fill in body markers.\n"
    "- body_starts_with: VERBATIM 30-80 chars from the document marking the start of the body.\n"
    "- body_ends_with: VERBATIM 30-80 chars marking the end.\n"
    "- These are used for substring search — copy exactly, character for character.\n\n"
    "STEP 4 — Fill organization vs author correctly.\n"
    "- organization = the institution (publication, board, court, agency).\n"
    "- author = the individual person (byline, signing official). Leave empty for "
    "institutional documents.\n\n"
    "Set is_news_content=false ONLY for files with zero journalistic value: "
    "empty files, binary garbage, personal shopping lists, error pages.\n"
    "Do NOT rewrite or summarize content — body text is sliced verbatim using your markers.\n"
    "You MUST call the register_stories tool. Do not respond with prose."
)


def _normalize_for_match(s: str) -> str:
    """Collapse whitespace so fuzzy substring searches survive small reformats."""
    return re.sub(r"\s+", " ", s).strip()


def _resolve_marker_offset(text: str, marker: str, *, after: int = 0) -> int:
    """Find `marker` in `text` starting at `after`, with whitespace-tolerant fallback.

    Returns -1 if not found.
    """
    marker = (marker or "").strip()
    if not marker:
        return -1

    # 1. Exact substring match.
    pos = text.find(marker, after)
    if pos >= 0:
        return pos

    # 2. Whitespace-normalized search — collapse runs of whitespace on both
    # sides and find the marker in the normalized text, then map the index
    # back to the original.
    norm_text = _normalize_for_match(text[after:])
    norm_marker = _normalize_for_match(marker)
    norm_pos = norm_text.find(norm_marker)
    if norm_pos >= 0:
        # Walk the original text counting normalized chars until we reach
        # norm_pos, then return the matching offset in the original.
        i, j = after, 0
        prev_space = False
        while i < len(text) and j < norm_pos:
            ch = text[i]
            if ch.isspace():
                if not prev_space:
                    j += 1
                    prev_space = True
            else:
                j += 1
                prev_space = False
            i += 1
        return i

    # 3. First-five-words fallback — handles minor LLM paraphrasing.
    first_words = " ".join(marker.split()[:5])
    if first_words and first_words != marker:
        pos = text.find(first_words, after)
        if pos >= 0:
            return pos

    return -1


def _slice_body(
    text: str,
    body_starts_with: str,
    body_ends_with: str,
    *,
    upper_bound: Optional[int] = None,
) -> str:
    """Locate the body inside `text` using the LLM's start/end snippets.

    `upper_bound`, when set, is the maximum end offset the body may extend
    to. This is how the caller prevents one story's body from bleeding into
    the next when the end marker fails to resolve cleanly — the caller knows
    where the next story starts and passes that as the ceiling.
    """
    start = _resolve_marker_offset(text, body_starts_with)
    if start < 0:
        return ""

    cap = upper_bound if upper_bound is not None else len(text)
    cap = min(cap, len(text))

    end_marker = (body_ends_with or "").strip()
    end_pos = _resolve_marker_offset(text, end_marker, after=start)
    if end_pos < 0:
        end = cap
    else:
        end = end_pos + len(end_marker)
        end = min(end, cap)

    if end <= start:
        return ""
    return text[start:end].strip()


def _make_chunks(text: str) -> list[str]:
    """Split text into chunks of at most WINDOW_SIZE chars.

    When the text carries record separators from a structured source
    (top-level JSON lists rendered by _extract_json), pack whole records
    into each chunk so a story body is never split across chunks — no
    overlap needed because records are pre-delimited.

    For unstructured text without separators, fall back to fixed-size
    overlapping windows so a body straddling the cut between window N
    and window N+1 is still fully contained in at least one of them.
    """
    if len(text) <= WINDOW_SIZE:
        return [text]

    if RECORD_SEPARATOR in text:
        records = text.split(RECORD_SEPARATOR)
        sep_len = len(RECORD_SEPARATOR)
        chunks: list[str] = []
        current: list[str] = []
        current_size = 0
        for rec in records:
            added_size = len(rec) + (sep_len if current else 0)
            if current and current_size + added_size > WINDOW_SIZE:
                chunks.append(RECORD_SEPARATOR.join(current))
                current = [rec]
                current_size = len(rec)
            else:
                current.append(rec)
                current_size += added_size
        if current:
            chunks.append(RECORD_SEPARATOR.join(current))
        return chunks

    chunks = []
    step = WINDOW_SIZE - WINDOW_OVERLAP
    start = 0
    while start < len(text):
        end = min(start + WINDOW_SIZE, len(text))
        chunks.append(text[start:end])
        if end >= len(text):
            break
        start += step
    return chunks


def _dedup_key(story: Story) -> str:
    """Stable key for deduplicating stories that show up in two overlapping
    chunks. Title + first 200 chars of body is specific enough to distinguish
    different events that happen to share a headline."""
    title = re.sub(r"\s+", " ", story.title.lower().strip())
    body_prefix = re.sub(r"\s+", " ", story.content[:200].lower().strip())
    return f"{title}||{body_prefix}"


def _normalize_chunk(
    text: str,
    source_label: str,
    anthropic_key: str,
    *,
    model: str = NORMALIZE_MODEL,
    link_hint: str = "",
    user_hint: str = "",
    allow_full_doc_fallback: bool = True,
) -> Tuple[list[Story], bool, str]:
    """Single LLM call on a chunk of text (which may be the whole document).

    allow_full_doc_fallback: when the LLM reports exactly one story but its
    markers fail to resolve, fall back to using this chunk's full text as
    that story's body. Safe only for whole-document calls; in chunked mode
    it would splice in content from adjacent stories.
    """
    client = chat_client(anthropic_key)

    user_prefix = f"Source label: {source_label}\n"
    if link_hint:
        user_prefix += f"Source URL: {link_hint}\n"
    if user_hint:
        user_prefix += f"Hint: {user_hint}\n"
    user_prefix += f"Text length (chars): {len(text)}\n\n"
    user_prefix += "----- BEGIN DOCUMENT -----\n"
    user_suffix = "\n----- END DOCUMENT -----"

    tool_use_block = None
    for attempt in range(NORMALIZE_MAX_ATTEMPTS):
        # Inner retry loop catches RateLimitError (Anthropic concurrent-
        # connection 429s) and logs every wait at WARNING so the operator
        # can see what's happening instead of staring at a silent spinner.
        resp = None
        for rl_attempt in range(RATE_LIMIT_MAX_RETRIES + 1):
            try:
                # Extended thinking is incompatible with forced tool_choice,
                # so we don't pass `thinking` here — omission == disabled.
                with ANTHROPIC_SEMAPHORE:
                    resp = client.messages.create(
                        model=model,
                        max_tokens=NORMALIZE_MAX_TOKENS,
                        system=[{
                            "type": "text",
                            "text": _NORMALIZE_SYSTEM,
                            "cache_control": {"type": "ephemeral"},
                        }],
                        messages=[
                            {"role": "user", "content": user_prefix + text + user_suffix},
                        ],
                        tools=[{**_NORMALIZE_TOOL, "cache_control": {"type": "ephemeral"}}],
                        tool_choice={"type": "tool", "name": "register_stories"},
                    )
                break
            except anthropic.RateLimitError as e:
                if rl_attempt >= RATE_LIMIT_MAX_RETRIES:
                    raise IngestError(
                        f"{source_label}: rate limit not cleared after "
                        f"{RATE_LIMIT_MAX_RETRIES} retries — {e}"
                    ) from e
                pause = rate_limit_pause(rl_attempt, e)
                logger.warning(
                    "Rate limited on %s; waiting %.0fs (attempt %d/%d).",
                    source_label, pause, rl_attempt + 1, RATE_LIMIT_MAX_RETRIES,
                )
                time.sleep(pause)
            except Exception as e:
                raise IngestError(
                    f"{source_label}: LLM normalization failed — {type(e).__name__}: {e}"
                ) from e

        assert resp is not None  # loop above either sets resp or raises

        tool_use_block = next(
            (b for b in resp.content if getattr(b, "type", None) == "tool_use"),
            None,
        )
        if tool_use_block is not None:
            break
        if attempt + 1 < NORMALIZE_MAX_ATTEMPTS:
            logger.warning(
                "Chunk for %s returned no tool call on attempt %d/%d; retrying.",
                source_label, attempt + 1, NORMALIZE_MAX_ATTEMPTS,
            )

    if tool_use_block is None:
        raise IngestError(
            f"{source_label}: LLM did not return structured stories after "
            f"{NORMALIZE_MAX_ATTEMPTS} attempts."
        )

    payload = tool_use_block.input if isinstance(tool_use_block.input, dict) else {}
    return _stories_from_payload(
        payload, text, source_label, link_hint, allow_full_doc_fallback
    )


def _stories_from_payload(
    payload: dict,
    text: str,
    source_label: str,
    link_hint: str,
    allow_full_doc_fallback: bool,
) -> Tuple[list[Story], bool, str]:
    """Convert a register_stories tool payload into Story objects.
    Shared by both the Anthropic and OpenAI normalization paths."""
    is_news = bool(payload.get("is_news_content", False))
    skip_reason = (payload.get("skip_reason") or "").strip()
    raw_stories = payload.get("stories") or []

    dropped = sum(1 for r in raw_stories if not isinstance(r, dict))
    if dropped:
        logger.warning("Dropped %d non-dict entries from %s's stories array.", dropped, source_label)
    raw_stories = [r for r in raw_stories if isinstance(r, dict)]

    if not raw_stories:
        return [], False, skip_reason or "No useful content found in this document."

    stories: list[Story] = []
    use_full_doc_fallback = allow_full_doc_fallback and len(raw_stories) == 1

    start_positions: list[int] = [
        _resolve_marker_offset(text, r.get("body_starts_with") or "")
        for r in raw_stories
    ]
    sorted_starts = sorted(s for s in start_positions if s >= 0)

    def _upper_bound_for(start: int) -> int:
        for s in sorted_starts:
            if s > start:
                return s
        return len(text)

    for idx, raw in enumerate(raw_stories):
        start = start_positions[idx]
        content = _slice_body(
            text,
            raw.get("body_starts_with") or "",
            raw.get("body_ends_with") or "",
            upper_bound=_upper_bound_for(start) if start >= 0 else None,
        )

        if len(content) < 20 and use_full_doc_fallback:
            logger.warning("Marker lookup for %s failed; falling back to full text.", source_label)
            content = text.strip()

        if not content:
            logger.warning("Dropping entry with empty body from %s.", source_label)
            continue

        title = (raw.get("title") or "").strip()
        if not title:
            first_line = next((ln for ln in content.splitlines() if ln.strip()), "")
            title = first_line[:80] if first_line else source_label

        ct = (raw.get("content_type") or "article").strip().lower()
        meta = raw.get("metadata")
        if not isinstance(meta, dict):
            meta = {}
        stories.append(Story(
            title=title,
            content=content,
            date=(raw.get("date") or "").strip(),
            author=(raw.get("author") or "").strip(),
            organization=(raw.get("organization") or "").strip(),
            link=(raw.get("link") or link_hint or "").strip(),
            content_type=ct if ct in CONTENT_TYPES else "other",
            metadata=meta,
            confidence=(raw.get("confidence") or "medium").strip(),
            reasoning=(raw.get("reasoning") or "").strip(),
        ))

    if not stories:
        return [], False, skip_reason or "No entry bodies could be extracted from this document."
    return stories, True, skip_reason




def normalize(
    text: str,
    source_label: str,
    anthropic_key: str,
    *,
    model: str = NORMALIZE_MODEL,
    concurrency: int = NORMALIZE_CONCURRENCY,
    link_hint: str = "",
    user_hint: str = "",
    on_progress: Optional[Callable[[dict], None]] = None,
) -> Tuple[list[Story], bool, str]:
    """Stage 2. Run extracted text through Haiku to produce structured entries.
    Single call for documents that fit in one window; for larger documents fan
    out into overlapping windows processed concurrently and deduplicate on merge.

    Returns (stories, is_news_content, skip_reason).
    """
    def _chunk_fn(chunk_text: str, fallback: bool) -> Tuple[list[Story], bool, str]:
        return _normalize_chunk(
            chunk_text, source_label, anthropic_key,
            model=model, link_hint=link_hint, user_hint=user_hint,
            allow_full_doc_fallback=fallback,
        )

    if not text or not text.strip():
        return [], False, "The document appears to be empty."

    if len(text) <= WINDOW_SIZE:
        if on_progress:
            on_progress({
                "stage": "normalize",
                "detail": "Processing 1 of 1 chunk",
                "completed": 1,
                "total": 1,
            })
        return _chunk_fn(text, True)

    chunks = _make_chunks(text)
    truncated = len(chunks) > MAX_CHUNKS
    if truncated:
        chunks = chunks[:MAX_CHUNKS]

    logger.info(
        "Chunked normalization for %s: %d chunks (text=%d chars, window=%d, overlap=%d)",
        source_label, len(chunks), len(text), WINDOW_SIZE, WINDOW_OVERLAP,
    )

    results: list[Optional[Tuple[list[Story], bool, str]]] = [None] * len(chunks)
    completed = 0
    completed_lock = threading.Lock()

    def _process(idx: int) -> Optional[Tuple[list[Story], bool, str]]:
        nonlocal completed
        try:
            return _chunk_fn(chunks[idx], False)
        except IngestError as e:
            logger.warning(
                "Chunk %d/%d failed for %s: %s",
                idx + 1, len(chunks), source_label, e,
            )
            return None
        finally:
            if on_progress:
                with completed_lock:
                    completed += 1
                    done = completed
                on_progress({
                    "stage": "normalize",
                    "detail": f"Processed {done} of {len(chunks)} chunks",
                    "completed": done,
                    "total": len(chunks),
                })

    with concurrent.futures.ThreadPoolExecutor(max_workers=concurrency) as ex:
        futures = {ex.submit(_process, i): i for i in range(len(chunks))}
        for fut in concurrent.futures.as_completed(futures):
            results[futures[fut]] = fut.result()

    failed = sum(1 for r in results if r is None)
    if failed == len(chunks):
        raise IngestError(
            f"{source_label}: all {len(chunks)} chunks failed during normalization."
        )

    # Merge: deduplicate by (title, body prefix). When the same story shows up
    # in two overlapping chunks, keep the one with the longer body — that's
    # the copy where the chunk window contained the full article rather than
    # clipping at the edge.
    merged: dict[str, Story] = {}
    skip_reasons: list[str] = []

    for result in results:
        if result is None:
            continue
        chunk_stories, _, chunk_skip = result
        if chunk_skip and not chunk_stories:
            skip_reasons.append(chunk_skip)
        for story in chunk_stories:
            key = _dedup_key(story)
            existing = merged.get(key)
            if existing is None or len(story.content) > len(existing.content):
                merged[key] = story

    stories = list(merged.values())

    if not stories:
        return [], False, (
            skip_reasons[0] if skip_reasons
            else "No news stories found in this document."
        )

    note_parts: list[str] = []
    if failed:
        note_parts.append(
            f"{failed} of {len(chunks)} chunks failed; partial results."
        )
    if truncated:
        note_parts.append(
            f"Document is very large; processed only the first {MAX_CHUNKS} chunks "
            f"(~{MAX_CHUNKS * (WINDOW_SIZE - WINDOW_OVERLAP) // 1000}k chars)."
        )

    return stories, True, " ".join(note_parts)


# ─────────────────────────────────────────────────────────────────────────────
# COMBINED API
# ─────────────────────────────────────────────────────────────────────────────


def ingest_file(
    filename: str,
    raw: bytes,
    anthropic_key: str,
    *,
    on_progress: Optional[Callable[[dict], None]] = None,
) -> IngestedSource:
    """Run both stages on an uploaded file. Never raises; failure is reported
    via excluded/skip_reason/extract_error on the returned IngestedSource."""
    source = IngestedSource(source_label=filename, kind="file")

    # Fast path: structured JSON with explicit story fields skips the LLM.
    if _ext_of(filename) == ".json":
        fast = _fast_json_stories(raw, filename)
        if fast is not None:
            source.stories = fast
            source.char_count = sum(len(s.content) for s in fast)
            if on_progress:
                on_progress({"stage": "done", "detail": "Structured JSON mapped without LLM"})
            return source

    # Build a filename-based type hint for PDFs and other documents so the LLM
    # can prioritize the right metadata schema from the start.
    is_pdf = _ext_of(filename) == ".pdf"
    user_hint = ""
    if is_pdf or _ext_of(filename) in (".doc", ".docx", ".rtf"):
        inferred = _infer_content_type_from_filename(filename)
        user_hint = f"Filename suggests this may be a '{inferred}'. Confirm and use the matching metadata fields."

    try:
        if on_progress:
            on_progress({"stage": "extract", "detail": "Extracting text"})
        text = extract_text(filename, raw)
    except IngestError as e:
        if is_pdf and "__SCANNED_PDF__" in str(e):
            # Fall back to OCR for scanned PDFs.
            if on_progress:
                on_progress({"stage": "extract", "detail": "Scanned PDF — running OCR"})
            try:
                text = _ocr_pdf(raw, anthropic_key, filename)
                if user_hint:
                    user_hint = "OCR transcription. " + user_hint
                else:
                    user_hint = "OCR transcription of a scanned PDF."
            except IngestError as ocr_err:
                source.excluded = True
                source.extract_error = str(ocr_err)
                source.skip_reason = str(ocr_err)
                return source
        else:
            source.excluded = True
            source.extract_error = str(e)
            source.skip_reason = str(e)
            return source

    source.char_count = len(text)
    source.truncated = len(text) > MAX_CHUNKS * (WINDOW_SIZE - WINDOW_OVERLAP)

    if not text.strip():
        source.excluded = True
        source.skip_reason = "The file appears to be empty or contains no readable text."
        return source

    try:
        stories, _, skip_reason = normalize(
            text, filename, anthropic_key,
            user_hint=user_hint,
            on_progress=on_progress,
        )
    except IngestError as e:
        source.excluded = True
        source.extract_error = str(e)
        source.skip_reason = "Normalization failed."
        return source

    source.stories = stories
    if not stories:
        source.excluded = True
        source.skip_reason = skip_reason or "No content entries could be extracted from this document."
    return source


def ingest_url(
    url: str,
    anthropic_key: str,
    *,
    on_progress: Optional[Callable[[dict], None]] = None,
) -> IngestedSource:
    """Fetch a URL, then run both stages. Same failure semantics as ingest_file."""
    source = IngestedSource(source_label=url, kind="url")

    try:
        filename, raw = fetch_url(url)
    except IngestError as e:
        source.excluded = True
        source.extract_error = str(e)
        source.skip_reason = "Failed to fetch this URL."
        return source

    try:
        if on_progress:
            on_progress({"stage": "extract", "detail": "Extracting text"})
        text = extract_text(filename, raw)
    except IngestError as e:
        source.excluded = True
        source.extract_error = str(e)
        source.skip_reason = str(e)
        return source

    source.char_count = len(text)
    source.truncated = len(text) > MAX_CHUNKS * (WINDOW_SIZE - WINDOW_OVERLAP)

    if not text.strip():
        source.excluded = True
        source.skip_reason = "The fetched page contains no readable text."
        return source

    try:
        stories, _, skip_reason = normalize(
            text, url, anthropic_key,
            link_hint=url,
            on_progress=on_progress,
        )
    except IngestError as e:
        source.excluded = True
        source.extract_error = str(e)
        source.skip_reason = "Normalization failed."
        return source

    source.stories = stories
    if not stories:
        source.excluded = True
        source.skip_reason = skip_reason or "No content entries could be extracted from this page."
    return source
